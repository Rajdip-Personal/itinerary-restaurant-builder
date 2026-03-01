#!/usr/bin/env python3
"""Generate a PowerPoint presentation of the Agentic AI Workshop flow diagram."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
NORDSTROM_DARK = RGBColor(0x1A, 0x1A, 0x2E)
HUMAN_BLUE = RGBColor(0x2D, 0x6A, 0x9F)
AGENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
HANDOFF_ORANGE = RGBColor(0xE6, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
SUBTLE_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def add_flow_card(slide, left, top, width, height, num, name, body, tag, color, bg_color, note=None):
    """Add a compact step card showing what the human does."""
    add_box(slide, left, top, width, height, bg_color, color, Pt(1.5))

    # Number badge
    bsz = Inches(0.28)
    badge = add_box(slide, left + Inches(0.07), top + Inches(0.07), bsz, bsz, color)
    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Step name
    add_text_box(slide, left + Inches(0.4), top + Inches(0.05), width - Inches(0.5), Inches(0.28),
                 name, font_size=11, bold=True, color=color)

    # Body text (what you do)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.36), width - Inches(0.2), Inches(0.4),
                 body, font_size=9, color=DARK_TEXT)

    # Note (e.g., "Skip if no UI") — above tag
    if note:
        ny = top + height - Inches(0.47) if tag else top + height - Inches(0.25)
        add_text_box(slide, left + Inches(0.1), ny, width - Inches(0.2), Inches(0.2),
                     note, font_size=8, color=SUBTLE_TEXT)

    # Tag (iterate/gate) — bottom of card
    if tag:
        add_text_box(slide, left + Inches(0.1), top + height - Inches(0.25), width - Inches(0.2), Inches(0.22),
                     tag, font_size=8, bold=True, color=color)


# ============================================================
# SLIDE 1: Your Workshop Journey (all 13 steps on one slide)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Title
add_text_box(slide, Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45),
             "Your Workshop Journey", font_size=28, bold=True, color=DARK_TEXT)
add_text_box(slide, Inches(0.3), Inches(0.5), Inches(12.7), Inches(0.28),
             "13 steps from PRD to working code \u2014 nothing moves forward without your approval",
             font_size=13, color=SUBTLE_TEXT)

# --- Phase 1: You + Claude ---
add_text_box(slide, Inches(0.3), Inches(0.9), Inches(6), Inches(0.28),
             "You + Claude  (Steps 1\u20135)", font_size=13, bold=True, color=HUMAN_BLUE)

card_h = Inches(1.2)
p1_w = Inches(2.25)
p1_gap = Inches(0.28)
p1_x0 = Inches(0.3)
p1_y = Inches(1.2)

phase1 = [
    (1, "Setup", "Verify environment\nPick team name", None, None),
    (2, "Pick Project", "Choose pre-built or\ncreate your own", None, None),
    (3, "Read PRD", "Everyone reads the PRD\nType \u2018ready\u2019 to continue", "\u23f8 GATE", None),
    (4, "Refine PRD", "Answer Claude\u2019s questions\nabout your PRD", "\u21bb Iterate until ready", None),
    (5, "Review Questions", "Resolve each open\nquestion one by one", "\u21bb Iterate per question", None),
]

for i, (num, name, body, tag, note) in enumerate(phase1):
    x = p1_x0 + i * (p1_w + p1_gap)
    add_flow_card(slide, x, p1_y, p1_w, card_h, num, name, body, tag, HUMAN_BLUE, LIGHT_BLUE_BG, note)
    if i < len(phase1) - 1:
        add_text_box(slide, x + p1_w, p1_y + Inches(0.35), p1_gap, Inches(0.3),
                     "\u2192", font_size=18, bold=True, color=HUMAN_BLUE, align=PP_ALIGN.CENTER)

# --- Handoff bar ---
hoff_y = Inches(2.55)
hbar = add_box(slide, Inches(0.3), hoff_y, Inches(12.733), Inches(0.38), LIGHT_ORANGE_BG, HANDOFF_ORANGE, Pt(2))
tf = hbar.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "HANDOFF \u2192  Agent Team coordinates from here.  Agents generate \u2014 you review and approve."
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = HANDOFF_ORANGE
p.alignment = PP_ALIGN.CENTER

# --- Phase 2: Agents generate, you approve ---
add_text_box(slide, Inches(0.3), Inches(3.05), Inches(8), Inches(0.28),
             "Agents Generate, You Approve  (Steps 6\u201313)", font_size=13, bold=True, color=AGENT_GREEN)

p2_w = Inches(2.95)
p2_gap = Inches(0.28)
p2_x0 = Inches(0.3)

# Row 1: Steps 6-9
p2_y1 = Inches(3.35)
row1 = [
    (6, "Execution Plan", "Review the plan\nApprove, revise, or re-run", "\u21bb Iterate", None),
    (7, "Requirements", "Review requirements\nApprove, revise, or re-run", "\u21bb Iterate", None),
    (8, "Technical Design", "Review the design\nApprove, revise, or re-run", "\u21bb Iterate", None),
    (9, "UI Prototype", "Click through prototype\nGive feedback", "\u21bb Iterate", "Skip if no UI"),
]

for i, (num, name, body, tag, note) in enumerate(row1):
    x = p2_x0 + i * (p2_w + p2_gap)
    add_flow_card(slide, x, p2_y1, p2_w, card_h, num, name, body, tag, AGENT_GREEN, LIGHT_GREEN_BG, note)
    if i < len(row1) - 1:
        add_text_box(slide, x + p2_w, p2_y1 + Inches(0.35), p2_gap, Inches(0.3),
                     "\u2192", font_size=18, bold=True, color=AGENT_GREEN, align=PP_ALIGN.CENTER)

# Row 2: Steps 10-13
p2_y2 = Inches(4.7)
row2 = [
    (10, "User Stories", "Review sprint-ready stories\nApprove, revise, or re-run", "\u21bb Iterate", None),
    (11, "Validation", "Review coverage report\nConfirm no gaps", "\u21bb Iterate", None),
    (12, "Jira Sync", "Provide Jira project key\nStories sync to Jira", None, None),
    (13, "Implementation", "Approve each story before\ncoding starts", "\u21bb Per story", None),
]

for i, (num, name, body, tag, note) in enumerate(row2):
    x = p2_x0 + i * (p2_w + p2_gap)
    add_flow_card(slide, x, p2_y2, p2_w, card_h, num, name, body, tag, AGENT_GREEN, LIGHT_GREEN_BG, note)
    if i < len(row2) - 1:
        add_text_box(slide, x + p2_w, p2_y2 + Inches(0.35), p2_gap, Inches(0.3),
                     "\u2192", font_size=18, bold=True, color=AGENT_GREEN, align=PP_ALIGN.CENTER)

# Legend
add_text_box(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
             "\u21bb = You can send it back for revision   \u2022   Blue = you talk directly to Claude   \u2022   Green = agents generate, you approve   \u2022   Every step commits & pushes to GitHub",
             font_size=10, color=SUBTLE_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Who Drives Each Step
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5),
             "Who Drives Each Step", font_size=28, bold=True, color=NORDSTROM_DARK)

add_text_box(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.35),
             "Rotate the keyboard. The person best suited to review each step should be driving.",
             font_size=14, color=SUBTLE_TEXT)

steps = [
    ("1", "Setup", "Engineer", HUMAN_BLUE),
    ("2", "Pick Project", "Whole team", HUMAN_BLUE),
    ("3", "Read PRD", "Everyone", HUMAN_BLUE),
    ("4", "Refine PRD", "Product Manager", HUMAN_BLUE),
    ("5", "Review Questions", "Product Manager", HUMAN_BLUE),
    ("6", "Execution Plan", "Tech Lead", AGENT_GREEN),
    ("7", "Requirements", "Tech Lead + PM", AGENT_GREEN),
    ("8", "Technical Design", "Architect / Sr Eng", AGENT_GREEN),
    ("9", "UI Prototype", "Whole team", AGENT_GREEN),
    ("10", "User Stories", "PM + Tech Lead", AGENT_GREEN),
    ("11", "Validation", "Tech Lead", AGENT_GREEN),
    ("12", "Jira Sync", "Anyone", AGENT_GREEN),
    ("13", "Implementation", "Engineer(s)", AGENT_GREEN),
]

y = Inches(1.2)
row_h = Inches(0.42)
for num, step, driver, color in steps:
    badge = add_box(slide, Inches(1.5), y, Inches(0.45), Inches(0.35), color)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.1), y, Inches(3), Inches(0.35),
                 step, font_size=13, bold=True, color=DARK_TEXT)

    add_text_box(slide, Inches(5.5), y, Inches(4), Inches(0.35),
                 driver, font_size=13, color=color, bold=True)

    phase = "You + Claude" if int(num) <= 5 else "Agents + You"
    phase_color = HUMAN_BLUE if int(num) <= 5 else AGENT_GREEN
    add_text_box(slide, Inches(9.5), y, Inches(2.5), Inches(0.35),
                 phase, font_size=10, color=phase_color)

    y += row_h


# ============================================================
# Save
# ============================================================
output_path = "docs/workshop-flow-diagram.pptx"
prs.save(output_path)
print(f"Created: {output_path}")
